"""Build the final-project presentation deck.

Run: python docs/report/build_pptx.py
Output: docs/report/presentation.pptx

Visual style — "Midnight Executive with Spotify accent":
  Primary:   #0F1419  (deep midnight, used for title/section/closing slides)
  Surface:   #F7F8FA  (light cream for content slides)
  Accent:    #1DB954  (Spotify green — callouts, key numbers, archetype dots)
  Warm:      #FF7A1A  (orange — secondary accent, the "user star" in mood scatter)
  Ink:       #1A1A1A  (body text on light)
  Muted:     #6B7280  (captions)
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parent
FIGS = ROOT / "figures"
OUT = ROOT / "presentation.pptx"

# ---- Palette ----
DARK = RGBColor(0x0F, 0x14, 0x19)
LIGHT = RGBColor(0xF7, 0xF8, 0xFA)
ACCENT = RGBColor(0x1D, 0xB9, 0x54)
WARM = RGBColor(0xFF, 0x7A, 0x1A)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x6B, 0x72, 0x80)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HEADER_FONT = "Calibri"
BODY_FONT = "Calibri"

# ---- Layout (16:9 slide is 13.333" × 7.5") ----
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def make_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation, bg: RGBColor):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.line.fill.background()
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    bg_shape.shadow.inherit = False
    return slide


def set_text(tf, text: str, *, size=18, bold=False, color=INK, font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def add_text(slide, text: str, x, y, w, h, *, size=18, bold=False, color=INK, font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    set_text(tb.text_frame, text, size=size, bold=bold, color=color, font=font, align=align, anchor=anchor)
    return tb


def add_lines(slide, lines, x, y, w, h, *, size=14, bold=False, color=INK, font=BODY_FONT, line_spacing=1.15, bullets=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        text = f"•  {line}" if bullets else line
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def accent_bar(slide, x, y, w, h, color=ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.shadow.inherit = False
    return s


def numbered_circle(slide, n: int, x, y, size=Inches(0.55), color=ACCENT, text_color=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.shadow.inherit = False
    set_text(
        s.text_frame, str(n), size=20, bold=True, color=text_color,
        font=HEADER_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    s.text_frame.margin_left = s.text_frame.margin_right = Inches(0)
    s.text_frame.margin_top = s.text_frame.margin_bottom = Inches(0)
    return s


def slide_title(slide, text: str, *, color=INK, subtitle: str | None = None, sub_color=MUTED):
    accent_bar(slide, Inches(0.6), Inches(0.65), Inches(0.12), Inches(0.55))
    add_text(
        slide, text, Inches(0.85), Inches(0.55), Inches(11.5), Inches(0.75),
        size=32, bold=True, color=color, font=HEADER_FONT, anchor=MSO_ANCHOR.MIDDLE,
    )
    if subtitle:
        add_text(
            slide, subtitle, Inches(0.85), Inches(1.2), Inches(11.5), Inches(0.4),
            size=14, color=sub_color, font=BODY_FONT, anchor=MSO_ANCHOR.TOP,
        )


def page_number(slide, n: int, total: int, on_dark=False):
    color = MUTED if not on_dark else RGBColor(0xA0, 0xA8, 0xB0)
    add_text(
        slide, f"{n} / {total}", Inches(12.6), Inches(7.05), Inches(0.65), Inches(0.3),
        size=10, color=color, font=BODY_FONT, align=PP_ALIGN.RIGHT,
    )


def footer_caption(slide, text: str, on_dark=False):
    color = MUTED if not on_dark else RGBColor(0xA0, 0xA8, 0xB0)
    add_text(
        slide, text, Inches(0.6), Inches(7.05), Inches(11.7), Inches(0.3),
        size=10, color=color, font=BODY_FONT,
    )


def add_card(slide, x, y, w, h, *, fill=CARD, accent: RGBColor | None = None):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    card.line.width = Pt(0.75)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.shadow.inherit = False
    if accent is not None:
        accent_bar(slide, x, y, Inches(0.08), h, accent)
    return card


def add_image(slide, path: Path, x, y, w=None, h=None):
    if not path.exists():
        raise FileNotFoundError(path)
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


# ============================================================================
TOTAL = 12


def slide_1_title(prs):
    s = blank_slide(prs, DARK)
    accent_bar(s, Inches(0.6), Inches(2.2), Inches(0.18), Inches(0.65), ACCENT)
    add_text(s, "Music Wrapped", Inches(0.85), Inches(2.05), Inches(11.5), Inches(0.95),
             size=54, bold=True, color=WHITE, font=HEADER_FONT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "An interactive music-discovery system with K-means personality archetypes",
             Inches(0.85), Inches(3.0), Inches(11.5), Inches(0.5),
             size=20, color=RGBColor(0xCA, 0xDC, 0xFC), font=BODY_FONT)
    add_text(s, "and cosine-similarity recommendations.",
             Inches(0.85), Inches(3.4), Inches(11.5), Inches(0.5),
             size=20, color=RGBColor(0xCA, 0xDC, 0xFC), font=BODY_FONT)

    add_text(s, "Jesse Boakye  ·  Kenneth Son  ·  Luis Zarate  ·  Ralph Lewis",
             Inches(0.85), Inches(5.6), Inches(11.5), Inches(0.4),
             size=15, bold=True, color=WHITE, font=BODY_FONT)
    add_text(s, "CPTS 315 — Data Mining   ·   Spring 2026   ·   Washington State University",
             Inches(0.85), Inches(6.0), Inches(11.5), Inches(0.4),
             size=12, color=RGBColor(0xA0, 0xA8, 0xB0), font=BODY_FONT)


def slide_2_problem(prs):
    s = blank_slide(prs, LIGHT)
    slide_title(s, "Beyond \"top 10 similar songs\"",
                subtitle="Why pair a content-based recommender with a Wrapped-style report?")

    # Left: text bullets
    add_lines(s, [
        "Most class-project recommenders stop at a ranked list of similar songs.",
        "Spotify Wrapped became a cultural moment because the presentation of audio analytics is engaging — not because the underlying recommender is novel.",
        "We pair a content-based recommender with a six-slide \"Wrapped\" report — mood, archetype, audio DNA, recs, hidden gems.",
        "Forces a complete data-mining pipeline: preprocessing, EDA, clustering, similarity, evaluation.",
    ], Inches(0.6), Inches(1.85), Inches(7.4), Inches(5.0), size=15, line_spacing=1.35, bullets=True)

    # Right: callout card with the question
    add_card(s, Inches(8.4), Inches(1.85), Inches(4.3), Inches(4.7), accent=ACCENT)
    add_text(s, "STATEMENT OF THE QUESTION", Inches(8.7), Inches(2.0), Inches(3.9), Inches(0.4),
             size=11, bold=True, color=ACCENT, font=BODY_FONT)
    add_text(s, "Can a content-based system using K-means archetypes and cosine-similarity ranking produce a recommendation experience that is:",
             Inches(8.7), Inches(2.45), Inches(3.9), Inches(1.6),
             size=13, color=INK, font=BODY_FONT)
    add_lines(s, [
        "algorithmically defensible — better than random on standard offline metrics?",
        "interpretable to a non-technical audience?",
        "worth shipping as an interactive demo?",
    ], Inches(8.7), Inches(4.05), Inches(3.9), Inches(2.4), size=12, line_spacing=1.25, bullets=True, color=INK)

    page_number(s, 2, TOTAL)


def slide_3_dataset(prs):
    s = blank_slide(prs, LIGHT)
    slide_title(s, "Dataset & technology stack",
                subtitle="Spotify Music Dataset (Kaggle) — static catalog, no API calls, content-based by necessity")

    # Three stat cards
    def stat(x, big, label):
        add_card(s, x, Inches(1.85), Inches(2.65), Inches(1.55), accent=ACCENT)
        add_text(s, big, x + Inches(0.2), Inches(1.95), Inches(2.45), Inches(0.85),
                 size=42, bold=True, color=ACCENT, font=HEADER_FONT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, label, x + Inches(0.2), Inches(2.85), Inches(2.45), Inches(0.5),
                 size=12, color=MUTED, font=BODY_FONT)

    stat(Inches(0.6), "4,494", "tracks after dedupe + clean")
    stat(Inches(3.4), "8", "shared audio features")
    stat(Inches(6.2), "26", "playlist-source genres")
    stat(Inches(9.0), "~20 KB", "persisted personality model")

    # Audio features used
    add_text(s, "MODELING FEATURES (MIN-MAX SCALED)", Inches(0.6), Inches(3.7), Inches(6), Inches(0.3),
             size=11, bold=True, color=ACCENT, font=BODY_FONT)
    add_text(s, "danceability  ·  energy  ·  valence  ·  acousticness  ·  instrumentalness  ·  speechiness  ·  tempo  ·  loudness",
             Inches(0.6), Inches(4.0), Inches(12.1), Inches(0.4),
             size=14, color=INK, font=BODY_FONT)

    # Tech stack
    add_text(s, "STACK", Inches(0.6), Inches(4.7), Inches(6), Inches(0.3),
             size=11, bold=True, color=ACCENT, font=BODY_FONT)
    add_lines(s, [
        "Python · pandas · NumPy · scikit-learn  —  preprocessing, K-means, PCA, cosine similarity, silhouette",
        "matplotlib · seaborn  —  EDA + report figures",
        "Streamlit · Plotly  —  interactive single-page demo with cached data + artifact loading",
        "joblib  —  persists scaler + KMeans + label map (`personality_v1.joblib`, ~20 KB)",
    ], Inches(0.6), Inches(5.0), Inches(12.1), Inches(2.0), size=13, line_spacing=1.3, bullets=True)

    page_number(s, 3, TOTAL)


def slide_4_pipeline(prs):
    s = blank_slide(prs, LIGHT)
    slide_title(s, "End-to-end pipeline",
                subtitle="Seven phases — every stage feeds a Wrapped slide")

    stages = [
        ("Preprocess", "Concat 2 CSVs · dedupe by track_id · MinMax scale 8 features"),
        ("EDA", "Distributions · correlations · genre balance · mood quadrants · PCA"),
        ("Cluster", "K-means K=2..10 · two tracks (raw vs PCA-6) · pick K=6"),
        ("Recommend", "User profile = mean(seed) · cosine sim · α=0.95 popularity blend"),
        ("Mood", "Rule-based valence/energy thresholds (0.6 high / 0.35 low)"),
        ("Hidden gems", "Top-decile similar · adaptive popularity ceiling · disjoint from recs"),
        ("Demo", "Streamlit single-page app · 6 sequential Wrapped slides"),
    ]
    y = Inches(1.95)
    for i, (name, desc) in enumerate(stages):
        row_y = y + Inches(0.66 * i)
        numbered_circle(s, i + 1, Inches(0.6), row_y, size=Inches(0.5))
        add_text(s, name, Inches(1.25), row_y - Inches(0.02), Inches(2.6), Inches(0.55),
                 size=17, bold=True, color=INK, font=HEADER_FONT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, desc, Inches(3.95), row_y - Inches(0.02), Inches(8.8), Inches(0.55),
                 size=13, color=MUTED, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    page_number(s, 4, TOTAL)


def slide_5_eda(prs):
    s = blank_slide(prs, LIGHT)
    slide_title(s, "EDA findings shaped the rest of the pipeline",
                subtitle="Three decisions came directly from looking at the data")

    # Left: correlation heatmap
    add_image(s, FIGS / "correlation_heatmap.png", Inches(0.55), Inches(1.85), w=Inches(5.6))
    add_text(s, "Audio feature correlations (8×8)", Inches(0.55), Inches(6.5), Inches(5.6), Inches(0.3),
             size=10, color=MUTED, font=BODY_FONT, align=PP_ALIGN.CENTER)

    # Right: three takeaways with accent dots
    items = [
        ("01", "Strong feature redundancy",
         "energy ↔ loudness ρ = 0.80 · acousticness ↔ energy ρ = -0.76. Six PCA components hit 95% variance — motivates the PCA-vs-original clustering experiment."),
        ("02", "Genre-skewed catalog",
         "Top 7 genres hold 300–560 tracks; sparse genres (disco, soca, mandopop) under 30. Bucketed into \"other\" for evaluation to avoid tiny-sample noise."),
        ("03", "Mood threshold tuning",
         "Valence × energy quadrants populated, but the textbook 0.5 split left low-energy/high-valence sparse. Retuned to high ≥ 0.6, low-valence < 0.35."),
    ]
    yt = Inches(1.85)
    for i, (n, head, body) in enumerate(items):
        row_y = yt + Inches(1.55 * i)
        # number badge
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.4), row_y, Inches(0.5), Inches(0.5))
        badge.line.fill.background()
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT
        badge.shadow.inherit = False
        set_text(badge.text_frame, n, size=12, bold=True, color=WHITE, font=HEADER_FONT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        badge.text_frame.margin_left = badge.text_frame.margin_right = Inches(0)
        add_text(s, head, Inches(7.0), row_y - Inches(0.02), Inches(5.7), Inches(0.4),
                 size=15, bold=True, color=INK, font=HEADER_FONT)
        add_text(s, body, Inches(7.0), row_y + Inches(0.42), Inches(5.7), Inches(1.05),
                 size=11.5, color=MUTED, font=BODY_FONT)

    page_number(s, 5, TOTAL)


def slide_6_cluster_pick(prs):
    s = blank_slide(prs, LIGHT)
    slide_title(s, "Choosing K — two-track experiment",
                subtitle="K-means on original 8 features (Track A) vs 6-component PCA (Track B), K = 2..10")

    add_image(s, FIGS / "kmeans_elbow_silhouette.png", Inches(0.5), Inches(1.85), w=Inches(7.5))
    add_text(s, "Elbow + silhouette across K. Highlighted band = demo target K ∈ [5, 8].",
             Inches(0.5), Inches(5.85), Inches(7.5), Inches(0.3),
             size=10, color=MUTED, font=BODY_FONT, align=PP_ALIGN.CENTER)

    # Right card: decision rule
    add_card(s, Inches(8.4), Inches(1.85), Inches(4.3), Inches(5.0), accent=ACCENT)
    add_text(s, "DECISION RULE", Inches(8.7), Inches(2.0), Inches(3.9), Inches(0.35),
             size=11, bold=True, color=ACCENT, font=BODY_FONT)
    add_lines(s, [
        "Pure silhouette peaks at K=2 — useless for archetype labels.",
        "Restrict to K ∈ [5, 8] (plan target).",
        "Drop any K with a cluster < 2% of catalog.",
        "Pick highest silhouette in the band.",
        "If two tracks within 0.02 silhouette: prefer Track A for interpretability.",
    ], Inches(8.7), Inches(2.45), Inches(3.9), Inches(2.6),
       size=11.5, color=INK, line_spacing=1.3, bullets=True)

    add_text(s, "WINNER", Inches(8.7), Inches(5.05), Inches(3.9), Inches(0.35),
             size=11, bold=True, color=ACCENT, font=BODY_FONT)
    add_text(s, "K = 6,  Track A", Inches(8.7), Inches(5.4), Inches(3.9), Inches(0.55),
             size=22, bold=True, color=INK, font=HEADER_FONT)
    add_text(s, "Silhouette 0.252 ± 0.0001 across 5 random seeds — stable.",
             Inches(8.7), Inches(6.0), Inches(3.9), Inches(0.7),
             size=11.5, color=MUTED, font=BODY_FONT)

    page_number(s, 6, TOTAL)


def slide_7_archetypes(prs):
    s = blank_slide(prs, DARK)
    accent_bar(s, Inches(0.6), Inches(0.55), Inches(0.18), Inches(0.65), ACCENT)
    add_text(s, "Six music personality archetypes",
             Inches(0.85), Inches(0.45), Inches(11.5), Inches(0.85),
             size=32, bold=True, color=WHITE, font=HEADER_FONT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "Manual labels driven by centroid feature profiles + per-cluster top-genre composition",
             Inches(0.85), Inches(1.15), Inches(11.5), Inches(0.4),
             size=13, color=RGBColor(0xCA, 0xDC, 0xFC), font=BODY_FONT)

    archetypes = [
        ("Lofi Wanderer", "high acoustic + instrumental, low energy. Lofi · jazz · world."),
        ("Night Driver", "low-acoustic, high-energy, low-valence. Moody mainstream."),
        ("Life of the Party", "highest danceability + valence + energy. Latin · pop · hip-hop."),
        ("The Old Soul", "very low energy, very high acoustic + instrumental. Wellness · classical · ambient."),
        ("The Sonic Architect", "electronic instrumentals · deep house. \"Here for the production, not the lyrics.\""),
        ("Acoustic Storyteller", "vocal acoustic mid-tempo. Latin · ambient · pop classics."),
    ]

    # 3 columns × 2 rows of cards on dark background
    card_w = Inches(4.05)
    card_h = Inches(2.55)
    gap_x = Inches(0.18)
    gap_y = Inches(0.18)
    total_w = card_w * 3 + gap_x * 2
    start_x = (prs.slide_width - total_w) // 2
    start_y = Inches(1.85)

    card_fill = RGBColor(0x1B, 0x22, 0x2A)  # slightly lighter than DARK
    for i, (name, desc) in enumerate(archetypes):
        col = i % 3
        row = i // 3
        x = start_x + (card_w + gap_x) * col
        y = start_y + (card_h + gap_y) * row
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, card_h)
        card.line.color.rgb = RGBColor(0x33, 0x3D, 0x47)
        card.line.width = Pt(0.5)
        card.fill.solid()
        card.fill.fore_color.rgb = card_fill
        card.shadow.inherit = False
        accent_bar(s, x, y, Inches(0.08), card_h, ACCENT)

        # number
        add_text(s, f"0{i+1}", x + Inches(0.3), y + Inches(0.2), Inches(1), Inches(0.4),
                 size=10, bold=True, color=ACCENT, font=BODY_FONT)
        # name
        add_text(s, name, x + Inches(0.3), y + Inches(0.55), card_w - Inches(0.45), Inches(0.6),
                 size=18, bold=True, color=WHITE, font=HEADER_FONT)
        # description
        add_text(s, desc, x + Inches(0.3), y + Inches(1.2), card_w - Inches(0.45), card_h - Inches(1.4),
                 size=12, color=RGBColor(0xCA, 0xDC, 0xFC), font=BODY_FONT)

    page_number(s, 7, TOTAL, on_dark=True)


def slide_8_recommendation(prs):
    s = blank_slide(prs, LIGHT)
    slide_title(s, "Recommendation engine",
                subtitle="Cosine similarity to the user's average profile, with a small popularity nudge")

    # Big formula card on left
    add_card(s, Inches(0.6), Inches(1.85), Inches(6.0), Inches(2.5), accent=ACCENT)
    add_text(s, "BLENDED SCORE", Inches(0.85), Inches(2.0), Inches(5.5), Inches(0.4),
             size=11, bold=True, color=ACCENT, font=BODY_FONT)
    add_text(s, "final(t) = α · cos(t, x̄)  +  (1-α) · pop(t)/100",
             Inches(0.85), Inches(2.45), Inches(5.5), Inches(0.65),
             size=20, bold=True, color=INK, font="Consolas", anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "x̄ = mean of seed audio features  ·  cos = cosine similarity  ·  pop = Spotify popularity in [0,100]",
             Inches(0.85), Inches(3.15), Inches(5.5), Inches(0.5),
             size=11, color=MUTED, font=BODY_FONT)
    add_text(s, "α = 0.95",
             Inches(0.85), Inches(3.6), Inches(5.5), Inches(0.6),
             size=22, bold=True, color=ACCENT, font=HEADER_FONT)

    # Right: alpha sweep image
    add_image(s, FIGS / "alpha_sweep.png", Inches(7.0), Inches(1.85), w=Inches(5.7))
    add_text(s, "Genre consistency vs α — focused users lose 47% relevance going 1.0 → 0.85.",
             Inches(7.0), Inches(5.55), Inches(5.7), Inches(0.3),
             size=10, color=MUTED, font=BODY_FONT, align=PP_ALIGN.CENTER)

    # Bottom: explainer
    add_text(s, "Why α = 0.95 (not the textbook 0.85):  alpha sweep showed α=0.85 cuts genre consistency from 0.16 to 0.08 for focused users — a 47% loss. α=0.95 keeps a small popularity nudge for cold-start without crushing relevance.",
             Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.95),
             size=12, color=INK, font=BODY_FONT, align=PP_ALIGN.LEFT)

    page_number(s, 8, TOTAL)


def slide_9_demo(prs):
    s = blank_slide(prs, LIGHT)
    slide_title(s, "Streamlit demo: six Wrapped slides",
                subtitle="`streamlit run app/streamlit_app.py` — single-file, single-page, sticky picker at top")

    slides = [
        ("Music in Numbers", "Counts, top genres, top artists from the user's selection."),
        ("Your Mood", "Mood label + valence/energy scatter against catalog backdrop."),
        ("Music Personality", "Archetype label + cluster centroid radar."),
        ("Audio DNA", "User profile vs. catalog mean — dual radar."),
        ("Top Recommendations", "Ranked table with similarity progress bars."),
        ("Hidden Gems", "Lower-popularity matches with similarity + popularity bars."),
    ]

    # 2 rows × 3 columns of cards
    card_w = Inches(4.0)
    card_h = Inches(2.45)
    gap_x = Inches(0.16)
    gap_y = Inches(0.2)
    total_w = card_w * 3 + gap_x * 2
    start_x = (prs.slide_width - total_w) // 2
    start_y = Inches(1.85)

    for i, (name, desc) in enumerate(slides):
        col = i % 3
        row = i // 3
        x = start_x + (card_w + gap_x) * col
        y = start_y + (card_h + gap_y) * row
        add_card(s, x, y, card_w, card_h, accent=ACCENT)
        add_text(s, f"SLIDE {i+1}", x + Inches(0.3), y + Inches(0.25), Inches(1.5), Inches(0.35),
                 size=10, bold=True, color=ACCENT, font=BODY_FONT)
        add_text(s, name, x + Inches(0.3), y + Inches(0.6), card_w - Inches(0.45), Inches(0.6),
                 size=18, bold=True, color=INK, font=HEADER_FONT)
        add_text(s, desc, x + Inches(0.3), y + Inches(1.25), card_w - Inches(0.45), card_h - Inches(1.4),
                 size=12, color=MUTED, font=BODY_FONT)

    add_text(s,
             "Catalog + personality artifact loaded once via @st.cache_resource · per-track genre hidden from rec/gems tables (dataset's playlist_genre is noisy)",
             Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.4),
             size=10, color=MUTED, font=BODY_FONT, align=PP_ALIGN.CENTER)

    page_number(s, 9, TOTAL)


def slide_10_evaluation(prs):
    s = blank_slide(prs, LIGHT)
    slide_title(s, "Evaluation — leave-one-out + metric comparison",
                subtitle="200 simulated users × 5 seeds = 1,000 LOO trials per strategy")

    add_image(s, FIGS / "loo_hitrate.png", Inches(0.55), Inches(1.85), w=Inches(6.2))
    add_text(s, "Hit Rate @ K (log scale) — genre-coherent vs random vs baseline.",
             Inches(0.55), Inches(5.55), Inches(6.2), Inches(0.3),
             size=10, color=MUTED, font=BODY_FONT, align=PP_ALIGN.CENTER)

    # Right: stat card
    add_card(s, Inches(7.1), Inches(1.85), Inches(5.6), Inches(5.0), accent=ACCENT)
    add_text(s, "HEADLINE NUMBERS  —  GENRE-COHERENT SEEDS",
             Inches(7.4), Inches(2.0), Inches(5.0), Inches(0.4),
             size=11, bold=True, color=ACCENT, font=BODY_FONT)

    metrics = [
        ("Hit @ 10", "1.6%", "8× random baseline (0.2%)"),
        ("Median rank", "950 / 4,494", "2.4× better than random (2,247)"),
        ("MRR", "0.0104", "vs ~0.001 for random seeds"),
        ("Recovered similarity", "0.936", "rank inflation = dataset density, not model failure"),
    ]
    yc = Inches(2.5)
    for i, (k, v, note) in enumerate(metrics):
        row_y = yc + Inches(1.05 * i)
        add_text(s, k, Inches(7.4), row_y, Inches(2.4), Inches(0.4),
                 size=12, bold=True, color=INK, font=BODY_FONT)
        add_text(s, v, Inches(7.4), row_y + Inches(0.35), Inches(2.4), Inches(0.55),
                 size=22, bold=True, color=ACCENT, font=HEADER_FONT)
        add_text(s, note, Inches(9.95), row_y + Inches(0.05), Inches(2.6), Inches(0.95),
                 size=11, color=MUTED, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    page_number(s, 10, TOTAL)


def slide_11_metric_compare(prs):
    s = blank_slide(prs, LIGHT)
    slide_title(s, "Cosine vs Euclidean — closing Phase 6",
                subtitle="Same 1,000 LOO trials, two ranking metrics")

    add_image(s, FIGS / "metric_comparison.png", Inches(0.55), Inches(1.85), w=Inches(6.5))
    add_text(s, "Hit Rate @ K (log scale) — cosine and Euclidean track closely.",
             Inches(0.55), Inches(5.7), Inches(6.5), Inches(0.3),
             size=10, color=MUTED, font=BODY_FONT, align=PP_ALIGN.CENTER)

    # Right: comparison table card
    add_card(s, Inches(7.4), Inches(1.85), Inches(5.3), Inches(5.0), accent=ACCENT)
    add_text(s, "WINNERS BY METRIC", Inches(7.7), Inches(2.0), Inches(4.7), Inches(0.4),
             size=11, bold=True, color=ACCENT, font=BODY_FONT)

    rows = [
        ("Hit @ 5",  "0.012",  "0.010", "cosine +20%"),
        ("Hit @ 10", "0.016",  "0.014", "cosine +14%"),
        ("MRR",      "0.0104", "0.0094","cosine +11%"),
        ("Hit @ 50", "0.066",  "0.072", "euclidean +9%"),
        ("Median rank", "951", "989",   "cosine"),
    ]
    yh = Inches(2.5)
    add_text(s, "metric",     Inches(7.7),  yh, Inches(1.4), Inches(0.3), size=11, bold=True, color=MUTED, font=BODY_FONT)
    add_text(s, "cosine",     Inches(9.05), yh, Inches(1.0), Inches(0.3), size=11, bold=True, color=ACCENT, font=BODY_FONT)
    add_text(s, "euclidean",  Inches(10.1), yh, Inches(1.05), Inches(0.3), size=11, bold=True, color=MUTED, font=BODY_FONT)
    add_text(s, "winner",     Inches(11.25), yh, Inches(1.4), Inches(0.3), size=11, bold=True, color=MUTED, font=BODY_FONT)
    for i, (k, c, e, w) in enumerate(rows):
        row_y = yh + Inches(0.5) + Inches(0.55 * i)
        add_text(s, k, Inches(7.7), row_y, Inches(1.4), Inches(0.4), size=12, color=INK, font=BODY_FONT)
        add_text(s, c, Inches(9.05), row_y, Inches(1.0), Inches(0.4), size=12, bold=True, color=ACCENT, font=BODY_FONT)
        add_text(s, e, Inches(10.1), row_y, Inches(1.05), Inches(0.4), size=12, color=INK, font=BODY_FONT)
        add_text(s, w, Inches(11.25), row_y, Inches(1.4), Inches(0.4), size=11, color=MUTED, font=BODY_FONT)

    add_text(s, "Decision: keep cosine. Precision wins where users actually look (top of the list), and the directional-similarity story is cleaner.",
             Inches(7.7), Inches(6.1), Inches(4.8), Inches(0.7),
             size=11.5, color=INK, font=BODY_FONT)

    page_number(s, 11, TOTAL)


def slide_12_close(prs):
    s = blank_slide(prs, DARK)
    accent_bar(s, Inches(0.6), Inches(0.55), Inches(0.18), Inches(0.65), ACCENT)
    add_text(s, "Limitations · future work · thank you",
             Inches(0.85), Inches(0.45), Inches(11.5), Inches(0.85),
             size=30, bold=True, color=WHITE, font=HEADER_FONT, anchor=MSO_ANCHOR.MIDDLE)

    # Two columns
    add_text(s, "WHAT'S HONEST ABOUT THE LIMITATIONS",
             Inches(0.6), Inches(1.5), Inches(6.0), Inches(0.4),
             size=11, bold=True, color=ACCENT, font=BODY_FONT)
    add_lines(s, [
        "No user behavior data — content-based by necessity.",
        "Dataset density: 4,494 tracks in [0,1]⁸ make absolute Hit@10 small (~1.6%).",
        "playlist_genre is noisy — the same song appears under different labels per source playlist. Hidden from per-track display.",
        "K-means and cosine treat all 8 features equally.",
    ], Inches(0.6), Inches(1.95), Inches(6.0), Inches(4.0),
       size=12, line_spacing=1.35, bullets=True, color=RGBColor(0xCA, 0xDC, 0xFC))

    add_text(s, "FUTURE WORK",
             Inches(7.0), Inches(1.5), Inches(5.7), Inches(0.4),
             size=11, bold=True, color=ACCENT, font=BODY_FONT)
    add_lines(s, [
        "Weighted cosine — variance-inverse over the seed selection.",
        "GMM personality — soft labels (\"60% Night Driver, 30% Lofi Wanderer\").",
        "Mood trajectory playlists — calm → upbeat sequencing.",
        "Popularity-residual hidden gems via regression.",
        "Live Spotify API integration with 30s previews.",
    ], Inches(7.0), Inches(1.95), Inches(5.7), Inches(4.0),
       size=12, line_spacing=1.35, bullets=True, color=RGBColor(0xCA, 0xDC, 0xFC))

    # Closing
    add_text(s, "Thank you  ·  Questions?",
             Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.7),
             size=32, bold=True, color=WHITE, font=HEADER_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "github.com/jesseboakye/315FinalProject",
             Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.4),
             size=12, color=ACCENT, font=BODY_FONT, align=PP_ALIGN.CENTER)


def main():
    prs = make_prs()
    slide_1_title(prs)
    slide_2_problem(prs)
    slide_3_dataset(prs)
    slide_4_pipeline(prs)
    slide_5_eda(prs)
    slide_6_cluster_pick(prs)
    slide_7_archetypes(prs)
    slide_8_recommendation(prs)
    slide_9_demo(prs)
    slide_10_evaluation(prs)
    slide_11_metric_compare(prs)
    slide_12_close(prs)

    prs.save(OUT)
    print(f"Saved -> {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
